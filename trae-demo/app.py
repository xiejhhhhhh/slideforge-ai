"""
SlideForge AI Demo Studio - 后端服务
提供 PPTX 模板解析、HTML 截图、PPTX 导出三个接口，并托管前端 index.html。

运行：python app.py
默认端口：5000
"""
from __future__ import annotations

import base64
import io
import os
import re
import tempfile
import traceback
import uuid
from pathlib import Path
from typing import Any

from flask import Flask, jsonify, request, send_from_directory, send_file
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------- 配置 ----------
BASE_DIR = Path(__file__).resolve().parent
OUTPUTS_DIR = BASE_DIR / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)

app = Flask(__name__, static_folder=None)
app.config["MAX_CONTENT_LENGTH"] = 256 * 1024 * 1024  # 256MB for real PPT templates


@app.after_request
def add_local_cors_headers(response):
    """Allow the local file:// demo page to call the Flask helper APIs."""
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization"
    return response


# ---------- 工具函数 ----------
def _hex_color(rgb) -> str:
    """python-pptx RGBColor -> #rrggbb"""
    try:
        return "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
    except Exception:
        return "#000000"


def _emu_to_pt(emu) -> float:
    """EMU -> pt"""
    try:
        return round(int(emu) / 12700, 1)
    except Exception:
        return 0.0


def _emu_to_inches(emu) -> float:
    try:
        return round(int(emu) / 914400, 2)
    except Exception:
        return 0.0


# ---------- 路由：前端静态服务 ----------
@app.route("/")
def index():
    return send_from_directory(BASE_DIR, "index.html")


@app.route("/<path:path>")
def static_files(path):
    return send_from_directory(BASE_DIR, path)


# ---------- 路由：健康检查 ----------
@app.route("/api/health")
def health():
    return jsonify({"ok": True, "service": "SlideForge AI Demo Studio"})


# ---------- 路由：解析 PPTX 模板 ----------
@app.route("/api/parse_template", methods=["POST"])
def parse_template():
    """解析上传的 .pptx 模板，提取比例/配色/字体/标题样式/布局/页脚等信息。"""
    if "file" not in request.files:
        return jsonify({"error": "未收到文件"}), 400
    f = request.files["file"]
    if not f.filename or not f.filename.lower().endswith(".pptx"):
        return jsonify({"error": "请上传 .pptx 文件"}), 400

    safe_suffix = Path(f.filename).suffix.lower() or ".pptx"
    tmp_path = OUTPUTS_DIR / "_uploaded_template_{}{}".format(uuid.uuid4().hex, safe_suffix)
    f.save(str(tmp_path))
    try:
        info = _analyze_pptx(tmp_path)
        info["template_path"] = str(tmp_path)
        return jsonify(info)
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "模板解析失败：{}".format(e)}), 500


def _analyze_pptx(path: Path) -> dict:
    """提取模板关键样式信息。"""
    prs = Presentation(str(path))
    info: dict[str, Any] = {
        "ratio": "16:9",
        "width_in": _emu_to_inches(prs.slide_width),
        "height_in": _emu_to_inches(prs.slide_height),
        "colors": [],
        "fonts": [],
        "title_style": {},
        "content_style": {},
        "footer": "",
        "suitable_pages": ["标题页", "内容页", "章节页", "结尾页"],
        "slide_count": len(prs.slides),
    }

    # 比例判断
    w, h = prs.slide_width, prs.slide_height
    ratio = round(w / h, 3) if h else 0
    if abs(ratio - 16/9) < 0.05:
        info["ratio"] = "16:9"
    elif abs(ratio - 4/3) < 0.05:
        info["ratio"] = "4:3"
    else:
        info["ratio"] = "{:.0f}:{:.0f}".format(w * 9, h * 16) if h else "未知"

    # 主题配色
    try:
        theme = prs.slide_masters[0].element
        # 提取主题色
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = theme.findall(".//a:clrScheme", ns)
        colors = []
        for cs in clr_scheme:
            for child in cs:
                tag = child.tag.split("}")[-1]
                if tag in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3"):
                    for c in child:
                        ctag = c.tag.split("}")[-1]
                        if ctag == "srgbClr":
                            colors.append("#" + c.get("val", "000000"))
                        elif ctag == "sysClr":
                            colors.append("#" + c.get("lastClr", "000000"))
        # 去重保留前 6 个
        seen = []
        for c in colors:
            if c not in seen:
                seen.append(c)
        info["colors"] = seen[:6]
    except Exception:
        info["colors"] = ["#1F2937", "#4F46E5"]

    # 遍历幻灯片提取字体/标题/内容样式
    fonts_set = []
    title_sizes = []
    content_sizes = []
    footer_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                for run in para.runs:
                    # 字体
                    try:
                        fnt = run.font
                        if fnt.name and fnt.name not in fonts_set:
                            fonts_set.append(fnt.name)
                        if fnt.size:
                            sz = _emu_to_pt(fnt.size)
                            if sz >= 24 and sz not in title_sizes:
                                title_sizes.append(sz)
                            elif 14 <= sz < 24 and sz not in content_sizes:
                                content_sizes.append(sz)
                    except Exception:
                        pass
                # 页脚识别
                low = text.lower()
                if any(k in low for k in ["page", "页", "©", "footer", "slide"]):
                    if len(text) < 40 and not footer_text:
                        footer_text = text

    info["fonts"] = fonts_set[:4] if fonts_set else ["Microsoft YaHei"]
    if title_sizes:
        info["title_style"] = {"font_size_pt": max(title_sizes), "font": info["fonts"][0]}
    else:
        info["title_style"] = {"font_size_pt": 32, "font": info["fonts"][0]}
    if content_sizes:
        info["content_style"] = {"font_size_pt": max(content_sizes), "font": info["fonts"][0]}
    else:
        info["content_style"] = {"font_size_pt": 18, "font": info["fonts"][0]}
    info["footer"] = footer_text

    return info


# ---------- 路由：HTML 截图 ----------
@app.route("/api/screenshot", methods=["POST"])
def screenshot():
    """把 HTML 字符串渲染为 PNG 截图（base64 返回）。"""
    data = request.get_json(silent=True) or {}
    html = data.get("html", "")
    width = int(data.get("width", 1280))
    height = int(data.get("height", 720))
    selector = data.get("selector", ".slide")

    if not html:
        return jsonify({"error": "html 为空"}), 400

    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return jsonify({"error": "服务器未安装 playwright，请运行 pip install playwright && playwright install chromium"}), 500

    tmp_html = OUTPUTS_DIR / "_screenshot_src.html"
    tmp_html.write_text(html, encoding="utf-8")

    try:
        images = []
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": width, "height": height}, device_scale_factor=2)
            page.goto(tmp_html.resolve().as_uri(), wait_until="networkidle")
            loc = page.locator(selector)
            count = loc.count()
            if count == 0:
                # 整页截图
                png_bytes = page.screenshot()
                images.append(base64.b64encode(png_bytes).decode("ascii"))
            else:
                for i in range(count):
                    png_bytes = loc.nth(i).screenshot()
                    images.append(base64.b64encode(png_bytes).decode("ascii"))
            browser.close()
        return jsonify({"images": images, "count": len(images)})
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "截图失败：{}".format(e)}), 500


# ---------- 路由：导出 PPTX ----------
@app.route("/api/export_pptx", methods=["POST"])
def export_pptx():
    """
    生成 PPTX 文件。
    入参 JSON:
      mode: "screenshot" | "native" | "template_fill"
      slides: [{ title, bullets[], notes, html?, time?, visual? }]
      slide_style: str (CSS)
      template_info: { width_in, height_in, colors[], fonts[], title_style, content_style, template_path? }
      images: [{ name, dataUrl, mime }]  # 用户上传的图片
      img_layout: [{ imgIndex, page, position, reason }]  # 图片布局建议
    """
    data = request.get_json(silent=True) or {}
    mode = data.get("mode", "screenshot")
    slides = data.get("slides", [])
    slide_style = data.get("slide_style", "")
    template_info = data.get("template_info", {}) or {}
    images = data.get("images", []) or []
    img_layout = data.get("img_layout", []) or []

    if not slides:
        return jsonify({"error": "没有幻灯片数据"}), 400

    try:
        if mode == "screenshot":
            pptx_path = _export_screenshot_mode(slides, template_info, slide_style, images, img_layout)
        elif mode == "template_fill":
            pptx_path = _export_template_fill_mode(slides, template_info, images, img_layout)
        else:
            pptx_path = _export_native_mode(slides, template_info, images, img_layout)

        return send_file(
            str(pptx_path),
            as_attachment=True,
            download_name="slideforge-demo.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "导出失败：{}".format(e)}), 500


def _get_slide_size(template_info: dict):
    """根据模板信息返回幻灯片尺寸（EMU）。"""
    w_in = template_info.get("width_in", 13.33)
    h_in = template_info.get("height_in", 7.5)
    return Inches(w_in), Inches(h_in)


def _decode_data_url(data_url: str) -> tuple[bytes, str]:
    """解析 data URL，返回 (bytes, mime)。"""
    # 格式：data:image/png;base64,xxxx
    if "," in data_url:
        header, b64 = data_url.split(",", 1)
        mime = "image/png"
        if ":" in header and ";" in header:
            mime = header.split(":")[1].split(";")[0]
        import base64 as _b64
        return _b64.b64decode(b64), mime
    raise ValueError("无效的 data URL")


def _save_images_to_tmp(images: list) -> list:
    """把前端传来的图片 dataUrl 保存为临时文件，返回 [{name, path, mime}]。"""
    saved = []
    for i, img in enumerate(images):
        try:
            data_url = img.get("dataUrl", "")
            if not data_url:
                continue
            data, mime = _decode_data_url(data_url)
            ext = ".png"
            if "jpeg" in mime or "jpg" in mime:
                ext = ".jpg"
            elif "svg" in mime:
                ext = ".svg"
            elif "gif" in mime:
                ext = ".gif"
            elif "webp" in mime:
                ext = ".webp"
            tmp = OUTPUTS_DIR / "_user_img_{:02d}_{}{}".format(i + 1, uuid.uuid4().hex[:8], ext)
            tmp.write_bytes(data)
            saved.append({
                "name": img.get("name", "image_{}".format(i + 1)),
                "path": str(tmp),
                "mime": mime,
                "index": i
            })
        except Exception as e:
            print("保存图片失败 {}: {}".format(img.get("name", "?"), e))
    return saved


def _parse_page_num(page_str: str) -> int:
    """从 '第3页' 或 '3' 提取页码，返回 1-based int。"""
    if not page_str:
        return 1
    m = re.search(r"\d+", str(page_str))
    if m:
        return int(m.group())
    return 1


def _position_to_rect(position: str, slide_w, slide_h, margin=0.4):
    """根据 position 描述返回 (left, top, width, height) EMU。"""
    pos = (position or "").lower()
    w = slide_w
    h = slide_h
    m = Inches(margin)
    # 默认右侧 40% 宽度
    if "右" in pos or "right" in pos:
        left = w * 0.58
        top = m
        width = w * 0.38
        height = h - 2 * m
    elif "左" in pos or "left" in pos:
        left = m
        top = m
        width = w * 0.38
        height = h - 2 * m
    elif "顶" in pos or "top" in pos or "上" in pos:
        left = m
        top = m
        width = w - 2 * m
        height = h * 0.35
    elif "底" in pos or "bottom" in pos or "下" in pos:
        left = m
        top = h * 0.6
        width = w - 2 * m
        height = h * 0.32
    elif "全屏" in pos or "背景" in pos or "full" in pos or "background" in pos:
        left = 0
        top = 0
        width = w
        height = h
    else:
        # 默认右侧
        left = w * 0.58
        top = m
        width = w * 0.38
        height = h - 2 * m
    return left, top, width, height


def _insert_images_for_slide(slide, page_num: int, saved_images: list, img_layout: list, slide_w, slide_h):
    """把指定页码的图片插入到 slide 中。"""
    for layout in img_layout:
        try:
            layout_page = _parse_page_num(layout.get("page", ""))
            if layout_page != page_num:
                continue
            img_idx = layout.get("imgIndex", 0)
            if img_idx >= len(saved_images):
                continue
            img_info = saved_images[img_idx]
            img_path = img_info["path"]
            position = layout.get("position", "右侧")
            left, top, width, height = _position_to_rect(position, slide_w, slide_h)
            # SVG 需要特殊处理，python-pptx 不直接支持 SVG，跳过或转 PNG
            if img_path.lower().endswith(".svg"):
                # 尝试用 playwright 转 PNG
                try:
                    png_path = _svg_to_png(img_path)
                    if png_path:
                        slide.shapes.add_picture(png_path, left, top, width=width, height=height)
                except Exception:
                    pass
            else:
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        except Exception as e:
            print("插入图片失败 (page={}): {}".format(page_num, e))


def _svg_to_png(svg_path: str) -> str | None:
    """把 SVG 转为 PNG（用 playwright）。"""
    try:
        from playwright.sync_api import sync_playwright
        svg_content = Path(svg_path).read_text(encoding="utf-8")
        html = '<!DOCTYPE html><html><head><meta charset="utf-8"><style>body{margin:0;padding:0;} svg{display:block;}</style></head><body>' + svg_content + '</body></html>'
        tmp_html = OUTPUTS_DIR / ("_svg_" + uuid.uuid4().hex + ".html")
        tmp_html.write_text(html, encoding="utf-8")
        png_path = str(Path(svg_path).with_suffix(".png"))
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1200, "height": 900}, device_scale_factor=2)
            page.goto(tmp_html.resolve().as_uri(), wait_until="networkidle")
            page.screenshot(path=png_path, full_page=True)
            browser.close()
        try:
            tmp_html.unlink()
        except Exception:
            pass
        return png_path
    except Exception as e:
        print("SVG 转 PNG 失败: {}".format(e))
        return None


def _export_screenshot_mode(slides: list, template_info: dict, slide_style: str = "", images: list = None, img_layout: list = None) -> Path:
    """截图模式：每页 HTML 截图为 PNG，嵌入 PPTX。HTML 中已包含 <img>，无需额外插图。"""
    from playwright.sync_api import sync_playwright

    slide_w, slide_h = _get_slide_size(template_info)
    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]

    # 收集所有页的 HTML
    full_html = '<!DOCTYPE html><html><head><meta charset="utf-8"><style>'
    full_html += "body{margin:0;background:#fff;}"
    full_html += ".slide{width:1280px;height:720px;box-sizing:border-box;overflow:hidden;}"
    full_html += slide_style or ""
    full_html += "</style></head><body>"
    for s in slides:
        full_html += s.get("html", "<div class='slide'></div>")
    full_html += "</body></html>"

    tmp_html = OUTPUTS_DIR / "_export_src.html"
    tmp_html.write_text(full_html, encoding="utf-8")

    img_paths = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        page.goto(tmp_html.resolve().as_uri(), wait_until="networkidle")
        loc = page.locator(".slide")
        count = loc.count()
        for i in range(count):
            img_path = OUTPUTS_DIR / "_export_slide_{:02d}.png".format(i + 1)
            loc.nth(i).screenshot(path=str(img_path))
            img_paths.append(img_path)
        browser.close()

    if not img_paths:
        raise RuntimeError("未截取到任何幻灯片")

    for img_path in img_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    out_path = OUTPUTS_DIR / "slideforge-demo.pptx"
    prs.save(str(out_path))
    return out_path


def _export_template_fill_mode(slides: list, template_info: dict, images: list = None, img_layout: list = None) -> Path:
    """模板填充模式：打开上传的 .pptx 模板，清空原 slides，用模板的布局添加新内容，严格套用模板风格，并插入用户图片。"""
    images = images or []
    img_layout = img_layout or []
    template_path_value = template_info.get("template_path")
    template_path = Path(template_path_value) if template_path_value else None
    if not template_path or not template_path.exists():
        candidates = sorted(OUTPUTS_DIR.glob("_uploaded_template_*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
        template_path = candidates[0] if candidates else OUTPUTS_DIR / "_uploaded_template.pptx"
    if not template_path.exists():
        raise RuntimeError("未找到上传的模板文件，请先在第 1 步上传 .pptx 模板")

    # 保存用户图片到临时文件
    saved_images = _save_images_to_tmp(images) if images else []

    # 打开模板作为基础（继承 master/layout/主题）
    prs = Presentation(str(template_path))

    # 清空模板原有的 slides（保留 master 和 layouts）
    xml_slides = prs.slides._sldIdLst
    for slide_id in list(xml_slides):
        xml_slides.remove(slide_id)

    # 选择可用 layout：优先用带标题+内容的布局（idx 1），否则用空白布局（idx 6），否则第 0 个
    layouts = prs.slide_layouts
    chosen_layout = layouts[0]
    for layout in layouts:
        ph_count = len(layout.placeholders)
        if ph_count >= 2:
            chosen_layout = layout
            break
    blank_layout = layouts[6] if len(layouts) > 6 else layouts[0]

    colors = template_info.get("colors") or ["#1F2937", "#4F46E5"]
    primary = colors[1] if len(colors) > 1 else colors[0]
    secondary = colors[2] if len(colors) > 2 else primary
    title_sz = (template_info.get("title_style") or {}).get("font_size_pt", 32)
    content_sz = (template_info.get("content_style") or {}).get("font_size_pt", 18)

    def _hex(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for idx, s in enumerate(slides):
        title = s.get("title", "幻灯片 {}".format(idx + 1))
        bullets = s.get("bullets", [])
        notes = s.get("notes", "")
        time_str = s.get("time", "")
        visual = s.get("visual", "")
        page_num = idx + 1

        # 第一页用标题页布局（idx 0 通常有标题占位符），其余用内容布局
        layout = chosen_layout if idx > 0 else (layouts[0] if layouts else chosen_layout)
        slide = prs.slides.add_slide(layout)

        # 填充占位符
        filled_title = False
        filled_body = False
        for shape in slide.placeholders:
            try:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0 and not filled_title:
                    shape.text = title
                    # 应用标题样式
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(title_sz)
                            run.font.bold = True
                            try:
                                run.font.color.rgb = _hex(primary)
                            except Exception:
                                pass
                    filled_title = True
                elif ph_idx == 1 and not filled_body and bullets:
                    tf = shape.text_frame
                    tf.clear()
                    for i, bullet in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(content_sz)
                        p.space_after = Pt(6)
                    filled_body = True
            except Exception:
                continue

        # 如果占位符没填上，用文本框兜底
        if not filled_title:
            txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(1.0))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(title_sz)
            p.font.bold = True
            try:
                p.font.color.rgb = _hex(primary)
            except Exception:
                pass
        if not filled_body and bullets:
            # 如果有图片要插入到这页，正文区缩小到左侧 55%
            has_img_this_page = any(_parse_page_num(l.get("page", "")) == page_num for l in img_layout)
            if has_img_this_page:
                body_left = Inches(0.6)
                body_top = Inches(1.6)
                body_width = slide_w * 0.5
                body_height = slide_h - Inches(2.4)
            else:
                body_left = Inches(0.6)
                body_top = Inches(1.6)
                body_width = prs.slide_width - Inches(1.2)
                body_height = prs.slide_height - Inches(2.4)
            txBox = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "• " + bullet
                p.font.size = Pt(content_sz)
                p.space_after = Pt(8)

        # 插入用户图片到这页（根据 img_layout）
        if saved_images and img_layout:
            _insert_images_for_slide(slide, page_num, saved_images, img_layout, slide_w, slide_h)

        # 视觉建议作为页脚提示（如果有）
        if visual:
            try:
                foot_left = Inches(0.6)
                foot_top = slide_h - Inches(0.45)
                foot_width = slide_w - Inches(1.2)
                foot_height = Inches(0.35)
                txBox = slide.shapes.add_textbox(foot_left, foot_top, foot_width, foot_height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "视觉：" + visual[:60]
                p.font.size = Pt(10)
                p.font.italic = True
                try:
                    p.font.color.rgb = _hex("#6B7280")
                except Exception:
                    pass
            except Exception:
                pass

        # 备注（包含时间和讲述重点）
        notes_full = notes or ""
        if time_str:
            notes_full = "【时间】" + time_str + "\n" + notes_full
        if notes_full:
            slide.notes_slide.notes_text_frame.text = notes_full

    out_path = OUTPUTS_DIR / "slideforge-demo.pptx"
    prs.save(str(out_path))
    return out_path


def _export_native_mode(slides: list, template_info: dict, images: list = None, img_layout: list = None) -> Path:
    """原生模式：用 python-pptx 绘制文本/形状，并插入用户图片。"""
    images = images or []
    img_layout = img_layout or []
    slide_w, slide_h = _get_slide_size(template_info)
    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]

    # 保存用户图片到临时文件
    saved_images = _save_images_to_tmp(images) if images else []

    colors = template_info.get("colors") or ["#1F2937", "#4F46E5"]
    primary = colors[1] if len(colors) > 1 else colors[0]
    title_font = (template_info.get("title_style") or {}).get("font", "Microsoft YaHei")
    title_sz = (template_info.get("title_style") or {}).get("font_size_pt", 32)
    content_font = (template_info.get("content_style") or {}).get("font", "Microsoft YaHei")
    content_sz = (template_info.get("content_style") or {}).get("font_size_pt", 18)

    def _hex(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        title = s.get("title", "幻灯片 {}".format(idx + 1))
        bullets = s.get("bullets", [])
        notes = s.get("notes", "")
        time_str = s.get("time", "")
        page_num = idx + 1

        # 检查这页是否有图片
        has_img_this_page = any(_parse_page_num(l.get("page", "")) == page_num for l in img_layout) if saved_images else False

        # 标题
        left = Inches(0.6)
        top = Inches(0.4)
        width = slide_w - Inches(1.2)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_sz)
        p.font.bold = True
        p.font.name = title_font
        p.font.color.rgb = _hex(primary)
        p.alignment = PP_ALIGN.LEFT

        # 内容（如果有图片，正文区缩小到左侧 55%）
        if has_img_this_page:
            body_left = Inches(0.6)
            body_top = Inches(1.8)
            body_width = slide_w * 0.5
            body_height = slide_h - Inches(2.6)
        else:
            body_left = Inches(0.6)
            body_top = Inches(1.8)
            body_width = slide_w - Inches(1.2)
            body_height = slide_h - Inches(2.6)
        txBox = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(content_sz)
            p.font.name = content_font
            p.font.color.rgb = _hex("#374151")
            p.space_after = Pt(8)
            p.alignment = PP_ALIGN.LEFT

        # 插入用户图片到这页（根据 img_layout）
        if saved_images and img_layout:
            _insert_images_for_slide(slide, page_num, saved_images, img_layout, slide_w, slide_h)

        # 页脚
        left = Inches(0.6)
        top = slide_h - Inches(0.5)
        width = slide_w - Inches(1.2)
        height = Inches(0.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        foot_text = "SlideForge AI · {}/{}".format(idx + 1, len(slides))
        if time_str:
            foot_text = "⏱ " + time_str + "  |  " + foot_text
        p.text = foot_text
        p.font.size = Pt(10)
        p.font.name = content_font
        p.font.color.rgb = _hex("#9CA3AF")
        p.alignment = PP_ALIGN.RIGHT

        # 备注
        notes_full = notes or ""
        if time_str:
            notes_full = "【时间】" + time_str + "\n" + notes_full
        if notes_full:
            slide.notes_slide.notes_text_frame.text = notes_full

    out_path = OUTPUTS_DIR / "slideforge-demo.pptx"
    prs.save(str(out_path))
    return out_path


# ---------- 入口 ----------
if __name__ == "__main__":
    print("=" * 50)
    print("SlideForge AI Demo Studio 后端服务")
    print("访问 http://localhost:5000 打开 Demo")
    print("按 Ctrl+C 退出")
    print("=" * 50)
    app.run(host="127.0.0.1", port=5000, debug=False)
