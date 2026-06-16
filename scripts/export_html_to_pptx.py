from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from render_html_screenshots import render


def export_pptx(html_path: Path, output_pptx: Path) -> None:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        image_dir = Path(tmp)
        render(html_path, image_dir)
        images = sorted(image_dir.glob("slide_*.png"))
        if not images:
            raise RuntimeError("No rendered slide images found.")

        prs = Presentation()
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for image in images:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)

        if len(prs.slides) > len(images):
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]

        prs.save(output_pptx)


def main() -> None:
    parser = argparse.ArgumentParser(description="Export an HTML slide deck to a PPTX file.")
    parser.add_argument("html", type=Path)
    parser.add_argument("output_pptx", type=Path)
    args = parser.parse_args()
    export_pptx(args.html, args.output_pptx)


if __name__ == "__main__":
    main()
